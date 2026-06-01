#!/usr/bin/env python
"""Assemble ordered slide images into a PPTX folder package."""

from __future__ import annotations

import argparse
import json
import re
import shutil
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}


def natural_key(path: Path) -> list[object]:
    parts = re.split(r"(\d+)", path.name.lower())
    return [int(part) if part.isdigit() else part for part in parts]


def collect_images(images_dir: Path) -> list[Path]:
    if not images_dir.exists() or not images_dir.is_dir():
        raise SystemExit(f"images-dir does not exist or is not a directory: {images_dir}")
    images = [
        path
        for path in images_dir.iterdir()
        if path.is_file() and path.suffix.lower() in IMAGE_EXTS
    ]
    images.sort(key=natural_key)
    if not images:
        raise SystemExit(f"No supported images found in {images_dir}")
    return images


def safe_name(name: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("._-")
    return cleaned or "slide"


def copy_ordered_images(images: list[Path], target_dir: Path) -> list[Path]:
    target_dir.mkdir(parents=True, exist_ok=True)
    copied: list[Path] = []
    for index, src in enumerate(images, 1):
        dst = target_dir / f"{index:03d}_{safe_name(src.stem)}{src.suffix.lower()}"
        shutil.copy2(src, dst)
        copied.append(dst)
    return copied


def add_image_slide(prs: Presentation, image_path: Path, fit: str) -> None:
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if fit == "stretch":
        slide.shapes.add_picture(str(image_path), 0, 0, width=slide_w, height=slide_h)
        return

    with Image.open(image_path) as img:
        img_w, img_h = img.size
    image_ratio = img_w / img_h
    slide_ratio = slide_w / slide_h

    if image_ratio >= slide_ratio:
        width = slide_w
        height = int(slide_w / image_ratio)
        left = 0
        top = int((slide_h - height) / 2)
    else:
        height = slide_h
        width = int(slide_h * image_ratio)
        left = int((slide_w - width) / 2)
        top = 0

    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def build_pptx(images: list[Path], pptx_path: Path, fit: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # Remove the default empty slide if the template ever creates one.
    while len(prs.slides) > 0:
        r_id = prs.slides._sldIdLst[0].rId  # python-pptx has no public remove API
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]

    for image in images:
        add_image_slide(prs, image, fit)

    pptx_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Create a PPTX package from ordered images.")
    parser.add_argument("--images-dir", required=True, help="Folder containing generated slide images.")
    parser.add_argument("--out-dir", required=True, help="Final output folder.")
    parser.add_argument("--deck-name", required=True, help="PPTX filename stem.")
    parser.add_argument("--image-subdir", default="images", help="Copied image subfolder name.")
    parser.add_argument("--fit", default="stretch", choices=["stretch", "contain"])
    parser.add_argument("--clean", action="store_true", help="Remove out-dir before packaging.")
    args = parser.parse_args()

    source_dir = Path(args.images_dir).resolve()
    out_dir = Path(args.out_dir).resolve()
    deck_name = safe_name(args.deck_name)

    if args.clean and out_dir.exists():
        shutil.rmtree(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    source_images = collect_images(source_dir)
    package_images_dir = out_dir / args.image_subdir
    if package_images_dir.exists():
        shutil.rmtree(package_images_dir)
    copied_images = copy_ordered_images(source_images, package_images_dir)

    pptx_path = out_dir / f"{deck_name}.pptx"
    build_pptx(copied_images, pptx_path, args.fit)

    manifest = {
        "deck_name": deck_name,
        "pptx": str(pptx_path),
        "images_dir": str(package_images_dir),
        "fit": args.fit,
        "slide_count": len(copied_images),
        "slides": [
            {
                "index": index,
                "source": str(src),
                "packaged": str(dst),
            }
            for index, (src, dst) in enumerate(zip(source_images, copied_images), 1)
        ],
    }
    manifest_path = out_dir / "manifest.json"
    manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")

    result = {
        "success": True,
        "output_folder": str(out_dir),
        "pptx": str(pptx_path),
        "images_dir": str(package_images_dir),
        "manifest": str(manifest_path),
        "slide_count": len(copied_images),
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
